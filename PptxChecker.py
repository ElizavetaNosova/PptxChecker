# -*- coding: utf-8 -*-
"""
Created on Sun Jun 14 18:43:35 2020

@author: lizan
"""

import tkinter as tk
from tkinter.filedialog import askopenfilename
import pptx
from nltk import sent_tokenize
import requests

from abc import ABC, abstractmethod
import re
from copy import copy
import os

class PrintOutputDecorator:
    '''
    Decorator, which prints the function's output before returning it.
    In this project this decorator is used to inform the user about his choice
    '''
    
    def __init__(self, message=''):
        self.message = message
        
    def __call__(self, foo):
        def wrapped_f(*args):
            foo_output = foo(*args)
            print(self.message)
            print(foo_output)
            return foo_output
        return wrapped_f

    
class AbstractCheckerRoot(ABC): 
    '''
    Organizes the checking process: asks the user to choose file
    and aspects, calls the checker, shows options, collects choices,
    corrects file and saves it.
    '''
    def __init__(self):
        self.checker_root = tk.Tk()
        self.checker_root.title('Проверка презентаций')
        info = tk.Label(self.checker_root, text = '''Последовательность действий:
            1. Выберите файл
            2. Выберите аспекты проверки
            3. Выберите варианты исправления либо откажитесь от исправления.
            Ошибки выводятся в том порядке, в каком встретились в тексте.
            Файл будет сохранён в той же папке и назван в формате
            "старое_название_corrected"''')
        info.pack()
        start_button = tk.Button(self.checker_root, 
                                 text='Приступить к выбору файла',
                                 command=self._start_button_click)
        start_button.pack()
        self.checker_root.mainloop()
        
    def _start_button_click(self):
        self.file_path = self._choose_file()
        self.chosen_file = self._open_file()
        self.prs_to_check = self._open_file()
        self._choose_options()
        
    def _choose_options(self):
        self._aspects_choice_window = tk.Toplevel(self.checker_root)
        self.__check_spelling = tk.BooleanVar(value=1)
        self.__check_typography = tk.BooleanVar(value=1)
        spelling_chk = tk.Checkbutton(self._aspects_choice_window,
                                      text='Проверка орфографии',
                                      variable=self.__check_spelling)
        spelling_chk.pack()
        typography_chk = tk.Checkbutton(self._aspects_choice_window,
                                        text='Проверка типографского оформления',
                                        variable=self.__check_typography)
        typography_chk.pack()
        start_button = tk.Button(self._aspects_choice_window,
                                 text='Приступить к проверке',
                                 command=self._start_checking)
        start_button.pack()
        self._aspects_choice_window.mainloop()
        
    def _start_checking(self):
        self._aspects_choice_window.destroy()
        checker_manager = AspectCheckerManager()
        chosen_aspects = self.__chosen_aspects()
        checker_manager.set_problems(self.chosen_file, chosen_aspects)
        self._show_options()
    
    @PrintOutputDecorator('Вы выбрали аспекты проверки:')
    def __chosen_aspects(self):
        chosen_aspects = []
        if self.__check_spelling.get():
            chosen_aspects.append('Проверка орфографии')
        if self.__check_typography.get():
            chosen_aspects.append('Проверка типографского оформления')
        return chosen_aspects
    
    @abstractmethod
    def _open_file(self):
        pass
    
    @abstractmethod
    def _show_options(self):
        pass


class PptxCheckerRoot(AbstractCheckerRoot):
    
    @PrintOutputDecorator('Вы выбрали файл:')
    def _choose_file(self):
        '''
        The method asks user to choose .pptx file
        '''
        filetypes=[('Презентация Microsoft PowerPoint', '.pptx')]
        title = 'Выберите презентацию в формате .pptx для проверки'
        return askopenfilename(filetypes=filetypes,
                               title = title)
    
    def _open_file(self):
        '''
        Creates PptxChecker
        '''
        return PptxChecker(self.file_path)
    
    def _show_options(self):
        '''
        Creates the window with shows problems and correcting options
        grouped by slides
        '''
        #self.checker_root.destroy()
        self.problems_window = PptxProblemsWindowCreator(self.checker_root,
                                                         self.chosen_file)
        
       
class SingleProblemAsker:
    '''
    The class places single problem data at the window.
    Stores the reference to the object with problem data.
    '''
    def __init__(self, problem, window, start_row):
        self.problem = problem
        
        self.text = tk.Label(window, text=self.problem['text_to_show'],
                             justify='left')
        self.text.grid(row=start_row)
        options = self.problem['options']
        self.var = tk.IntVar()
        self.var.set(0)
        for i, option in enumerate(options):
            radiobutton = tk.Radiobutton(window, text=options[i],
                                         variable = self.var, value=i)
            if i == 0:
                radiobutton.select()
            radiobutton.grid(row=start_row+1, column=i)
        
    def add_choice(self):
        '''
        The method interpretes Tkinter variables and adds key 'choice'
        to the corresponding dictionary
        '''
        try:
            self.problem['choice'] = self.var.get()
        except ProtectedKeyException:
            pass       


class AbsractProblemsWindowCreator():
    
    '''
    The class organizes showing the problems and options to the user
    and collecting his choice data
    '''
    def __init__(self, root, prepared_object):
        self.window = tk.Toplevel(root)
        self.start_row = 0
        self.single_askers = []
        self.__correction_started = False
        self.collect_user_choises(prepared_object)
        self.window.mainloop()
        
        
    def _show_problems(self, found_problems):
        '''
        Input - list of dictionaries.
        Method creates a SingleProblemAsker for each problem
        '''
        for problem in found_problems:
            problem_asker = SingleProblemAsker(problem, self.window,
                                               self.start_row)
            self.single_askers.append(problem_asker)
            self.start_row += 2
            
                
    @abstractmethod       
    def _show_file_problems(self, prepared_object):
        '''
        Shows all the problems and comments about its location
        in corresponding file
        '''
        pass
    
    def collect_user_choises(self, prepared_object):
        '''
        Organizes all the process of collecting user choices
        '''
        self.prepared_object = prepared_object
        self._show_file_problems(prepared_object)
        choice_button = tk.Button(self.window, text = 'Внести исправления',
                                  command=self.__prepare_data_to_correct)
        choice_button.grid(row=self.start_row)
    
    def __prepare_data_to_correct(self):
        '''
        For each single problem asker calls the method which adds
        the user choice data to the problem dict
        '''
        for asker in self.single_askers:
            asker.add_choice()
        self.__correct()
            
    def __correct(self):
        '''
        Calls the prepared object's method '__correct'
        '''
        if not self.__correction_started:
            self.__correction_started = True
            self.prepared_object.correct()
        self.window.destroy()
        
            
class PptxProblemsWindowCreator(AbsractProblemsWindowCreator):
    
    def _show_file_problems(self, prepared_pptx):
        '''
        Shows problems data and comments about the presentation
        slides numbers
        '''
        prev_slide_id = -1
        problems = prepared_pptx.problems_to_show()
        for index, shape_problems in enumerate(problems):
            current_slide_id = shape_problems[1]
            if current_slide_id > prev_slide_id:
                text = 'СЛАЙД '+str(current_slide_id+1)
                bondary_label = tk.Label(self.window, text = text)
                bondary_label.grid(row=self.start_row)
                self.start_row +=1
                prev_slide_id = current_slide_id
            if shape_problems[0]:
                shape_found_problems = shape_problems[0]
                self._show_problems(shape_found_problems)
            
            
class AbstractFileChecker(ABC):
    '''
    The class reads the file, extracts the texts, 
    converts it into the format, which is sutable for the
    aspect checkers (methods texts_to_checker, header_to_checker),
    stores checkers output(set_texts_problems, set_headers_problems)
    '''
    
    def __init__(self, file_path):
        self.file_path = file_path
        self._content = self._read_file(self.file_path)
        self._get_all_texts()
        self.__is_corrected = False
    
    @abstractmethod
    def _read_file(self, file_path):
        '''
        Reads the file
        '''
        pass
    
    @abstractmethod
    def texts_to_checker(self):
        '''
        Makes copies of text data in self.__texts and converts it
        into a list of strings
        '''
        pass   
     
    @abstractmethod    
    def set_texts_problems(self, texts_problems):
        '''
        Transforms the checker output into appropriate format and
        set is as self.__headers_problem
        '''
        pass
    
    @abstractmethod
    def problems_to_show(self):
        '''
        Joins all the problems and sorts is.
        '''
        pass
    
    def __correct_single_problem(self, text, problem): 
        '''
        Substitude mistakes in the text with chosen correction option.
        '''
        if problem['choice'] == len(problem['s']): ##The last option is always
            pass                                   ## 'do not correct'
        else:
            if (problem['pos'] + problem['len']) < len(text):
                post_part = text[problem['pos']+problem['len']:] 
            else: ###If the mistake is in the last symbol
                post_part = '' ##the post_part should be ''
            text = text[:problem['pos']] + problem['s'][problem['choice']] + post_part
        return(text)
    
    def _correct_single_text_problems(self, text, problems):
        '''
        Sorts the problems by pos decrising and corrects it.
        No overlaping of fragments to substitute should be guatanteed.
        '''
        problems.sort(key = lambda x: x['pos'], reverse=True)
        for problem in problems:
            text = self.__correct_single_problem(text, problem)
        return text
        
    def correct(self):
        '''
        Corrects content if it has not been corrected yet and saves it
        '''
        if self.__is_corrected == False:
            self._correct_content()
            self.__is_corrected = True
            self._save_content(self.__get_default_new_file_name())
        
        
    @abstractmethod   
    def _correct_content(self):
        '''
        Corrects all the texts and modifies self._content
        '''
        pass
    
    @abstractmethod
    def _save_content(self, new_file_name):
        pass
            
    def __get_default_new_file_name(self):
        '''
        Returns the name which will be given to the corrected file
        '''
        name_parts = os.path.splitext(self.file_path)
        return name_parts[0] + '_corrected'+ name_parts[1]
        

class PptxChecker(AbstractFileChecker):
    
    def _read_file(self, file_path):
        with open(file_path, 'rb') as f:
            return pptx.Presentation(f)

    def _get_all_texts(self):
        texts = []
        shapes = []
        for index, slide in enumerate(self._content.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    texts.append([shape.text, index])
                    shapes.append(shape)
        self.__shapes = shapes
        self.__texts = texts
        
    def texts_to_checker(self):
        return [copy(text[0]) for text in self.__texts]
    
    def set_texts_problems(self, texts_problems):
        self.__texts_problems = [[texts_problems[i], self.__texts[i][1]]
        for i, text in enumerate(self.__texts)]     
        
    def problems_to_show(self):
        return self.__texts_problems  
    
    def _correct_content(self):  
        for index, text in enumerate(self.__texts):
            text_problems = self.__texts_problems[index][0]
            if text_problems:
                new_text = self._correct_single_text_problems(text[0],
                                                              text_problems)
                self.__shapes[index].text = new_text
            
    def _save_content(self, new_file_name):
        self._content.save(new_file_name)
        

class AbstractAspectChecker(ABC):
    '''
    Forms dictionaries with text problems data
    '''
    
    def find_problems(self, texts):
        '''
        Input: a list of texts
        Output: a list of lists dictionaries, each list corresponds to one text,
        each dictionaty contains information about 1 problem
        '''
        problems = self._check_texts(texts)
        for i, single_text_problems in enumerate(problems):
            text = texts[i]
            for problem in single_text_problems:
                problem['options'] = self._get_options(problem)
                problem['text_to_show'] = self._get_text_to_show(text, problem)
        return problems
    
    @abstractmethod
    def _check_texts(self, texts):
        '''
        Finds the problems, its position ('pos'), length('len'),
        possible stings to substitute('s') and optionally some specific
        data for this type of problems
        '''
        pass
    
    def _get_options(self, problem):
        '''
        Get options: options to choose the replacements ('s')
        and extra option 'do not correct'
        '''
        options = copy(problem['s'])
        options.append('Не исправлять')
        return options
    
    
    def _get_sentences_with_id(self, text):
        '''
        Splits text into shorter fragments
        '''
        if len(text) <= 70:
            return [{'sent': text, 'pos':0, 'end': len(text)-1}]
        else:
            sents = sent_tokenize(text, 'russian')
            return [{'sent': sent, 'pos': text.find(sent),
                     'end': text.find(sent)+len(sent)-1} for sent in sents]
               
    def _get_text_to_show(self, text, problem):
        '''
        Matches a problem with a text fragment
        '''
        sents_with_index = self._get_sentences_with_id(text)
        current_sent_id = 0
        current_sent = sents_with_index[current_sent_id]
        while problem['pos'] > current_sent['end']:
            current_sent_id += 1
            current_sent = sents_with_index[current_sent_id]
        return current_sent['sent']
    

class ProtectedKeyException(Exception):
    def __init__(self, protected_key):
        self.text = 'Key '+protected_key+' cannot be overwrited'

### Вообще-то без этого класса можно было бы легко обойтись (есть защита 
### от повторного вызова исправления), но мне нужно было что-нибудь переопределить
###
        
class DictWithProtectedKey(dict):
    '''
    Dict with a key, which can not be overwrited.
    In this project it is used to be sure that system do not ask about the same
    problem twice
    '''
    def __init__(self, data, protected_key='choice'):
        super().__init__(data)
        self.__protected_key = protected_key
        self.__protecked_key_is_used = True if protected_key in self else False
        
    def __setitem__(self, key, value):
        if key != self.__protected_key:
            super().__setitem__(key, value)
        else:
            if not self.__protecked_key_is_used:
                self.__protecked_key_is_used = True
                super().__setitem__(key, value)
            else:
                raise(ProtectedKeyException)
    
    def get_protected_key(self):
        return copy(self.__protected_key)
                                        

class SpellChecker(AbstractAspectChecker):  
    '''
    In  this project YandexSpeller is used.
    Output of other checkers has the same type
    '''

    def _check_texts(self, texts, CHECKER_LIMIT=10000):
        '''
        Makes a query to YandexSpeller, if the texts do not exceed its limit.
        It they do, join the output of several queries
        '''
        if sum([len(text) for text in texts]) < CHECKER_LIMIT:
            problems = self.__checker_query(texts)
        else:
            problems = self.__split_and_check(texts)
        for text_problems in problems:
            for problem in text_problems:
                problem['type'] = 'spelling'
        return problems
    
    def __checker_query(self, texts):
        response = requests.post('https://speller.yandex.net/services/spellservice.json/checkTexts',
                                 {'text': texts}).json()
        problems = [[DictWithProtectedKey(problem) for problem in text_problems] 
        for text_problems in response]
        return problems
    
    def __split_and_check(self, texts, CHECKER_LIMIT):
        if all([len(text) for text in texts]) < CHECKER_LIMIT:
            responses = []
            current_texts = []
            current_len = 0
            for i, text in enumerate(texts):
                if current_len + len(text) > CHECKER_LIMIT:
                    responses += self.__checker_query(current_texts)
                    current_len = 0
                    current_texts = []
                current_texts.append(text)
            responses += self.__checker_query(current_texts)
            return responses
        else:
            raise Exception('В тексте не должно быть блоков длиннее '+
                            str(CHECKER_LIMIT) + ' символов')
            
    def _get_text_to_show(self, text, problem):
        text_fragment = super()._get_text_to_show(text, problem)
        word_note = 'Опечатка в слове ' + problem['word'] + ': '
        return word_note + text_fragment      

class TypographyChecker(AbstractAspectChecker):
    '''
    Checks type of dash and quotation.
    Do not check punctuation (only typography)
    '''
       
    def __init__(self):
        self.__problem_comments = {'hypher_between_spaces':
            'Дефис между пробелами. Рекомендуем заменить дефис на тире либо убрать пробелы.',
                                   'dash_between_digits':
                                       'Между цифрами ставится среднее тире (значение от ... до ...) либо дефис (приблизительное значение)',
                                   'quotation_type': 'Неверный тип кавычек.'}
    
    
    def _check_texts(self, texts):
        return [(self.__check_dash(text) + self.__check_quotation(text))
                for text in texts]
    

    def __check_dash(self, text):
        dash_problems = []
        
        hypher_between_spaces = self.__find_all_matches(' - ', text)
        for problem in hypher_between_spaces:
            problem['s'] = [' — ', '-']
            problem['type'] = 'hypher_between_spaces'
        dash_problems += hypher_between_spaces
        
        dash_between_digits = self.__find_all_matches('\d[-—]\d', text)
        for problem in dash_between_digits:
            problem['s']  = ['–', '-']
            problem['len'] = 1
            problem['pos'] += 1
            problem['type'] = 'dash_between_digits'
        dash_problems += dash_between_digits
        
        return dash_problems
    
    def __check_quotation(self, text): 
        quotation_problems = []
        general_state = 0
        inner_state = 0
        wrong_inner_state = 0
        for i, symbol in enumerate(text):
            if symbol in '„“"':
                if general_state==0:
                    problem = DictWithProtectedKey({'pos': i, 'type': 'quotation_type', 'len': 1})
                    if wrong_inner_state==0:
                        problem['s'] = ['«', '']
                        wrong_inner_state = 1
                    else:
                        problem['s'] = ['»', '']
                        wrong_inner_state = 0
                    quotation_problems.append(problem)
                else:
                    if inner_state==0:
                        inner_state=1
                        if symbol == '“':
                            problem = DictWithProtectedKey({'pos': i, 'type': 'quotation_type',
                                       'len': 1, 's': ['„', '']})
                            quotation_problems.append(problem)
                    else:
                        inner_state=0
                        if symbol == '„':
                            problem = DictWithProtectedKey({'pos': i, 'type': 'quotation_type',
                                       'len': 1, 's': ['“', '']})
                            quotation_problems.append(problem)
            elif symbol == '«':
                if general_state == 1:
                    problem = DictWithProtectedKey({'pos': i, 'type': 'quotation_type', 'len': 1,
                               's': ['„', '»', '']})
                    quotation_problems.append(problem)
                general_state = 1
            elif symbol == '»':
                if general_state == 0:
                    problem = DictWithProtectedKey({'pos': i, 'type': 'quotation_type', 'len': 1,
                               's': ['«', '']})
                    quotation_problems.append(problem)
                general_state = 0
        return quotation_problems
    
        
    def _get_text_to_show(self, text, problem):
        text_to_show = super()._get_text_to_show(text, problem)
        text_to_show = self.__problem_comments[problem['type']] +'\n' + text_to_show
        return text_to_show
                
    def _get_options(self, problem):
        if problem['type'] == 'quotation_type':
            options = problem['s'][:-1] + ['Удалить кавычку', 'Не исправлять']
        else:
            options = super()._get_options(problem)
        return options
        
            
    def __find_all_matches(self, pattern, string):
        '''
        Returns dicts with positiion and length of all the matches
        '''
        matches = []
        current_match = re.search(pattern, string)
        while current_match:
            matches.append({'pos': current_match.start(),
                            'len': current_match.end()-current_match.start()})
            if current_match.end() < len(string):
                current_match = re.search(pattern,
                                          string[current_match.end():])
            else:
                current_match = None
        return matches
    

class AspectCheckerManager:
    '''
    Joins output of SpellChecker() and TypographyChecker(). 
    '''
    
    def __init__(self):
        self.spell_checker = SpellChecker()
        self.typography_checker = TypographyChecker()
    
    def set_problems(self, file_to_check, chosen_aspects):
        
        texts = file_to_check.texts_to_checker() 
        spelling_problems = self.__check_aspect(texts, self.spell_checker,
                                                ('Проверка орфографии' in chosen_aspects))
        typography_problems = self.__check_aspect(texts, self.typography_checker,
                                                  ('Проверка типографского оформления' in chosen_aspects)) 
        texts_problems = self.__sum_problems(spelling_problems,
                                             typography_problems)
        file_to_check.set_texts_problems(texts_problems)
    
    def __check_aspect(self, texts, checker, var):
        if var:
            return checker.find_problems(texts)
        else:
            return [[] for text in texts]
        
    def __sum_problems(self, problems1, problems2):
        problems = [self.__sorted_problems(problems1[i]+problems2[i])
        for i in range(len(problems1))] ##Length of the input lists 
        return problems                 ##are equal
    
    def __sorted_problems(self, problems): 
        return sorted(problems, key=lambda x: x['pos'])
    
    
if __name__ == '__main__':
    root = PptxCheckerRoot()
